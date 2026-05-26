"""Generate Session 3 closeout presentation.

Produces a ~12-slide deck summarizing what shipped in CRISP-DM Phase 3 for the
Mycorrhizal Barcelona project, including the PRPI v1.1 → v1.2 pivot.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "session-3-presentation.pptx"

# Brand palette — earthy / mycorrhizal
INK = RGBColor(0x1A, 0x2A, 0x1F)         # near-black green
MOSS = RGBColor(0x3D, 0x6B, 0x3A)        # primary green
AMBER = RGBColor(0xC8, 0x7E, 0x2E)       # accent
BONE = RGBColor(0xF5, 0xF1, 0xE8)        # warm off-white
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
RED = RGBColor(0xB2, 0x3A, 0x2B)


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def add_bg(slide, color=BONE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_accent_bar(slide, color=MOSS, x=0, y=0, w=0.2, h=7.5):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text(slide, text, x, y, w, h, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(slide, items, x, y, w, h, *, size=14, color=INK, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_footer(slide, page_num, total):
    add_text(slide, "Mycorrhizal Barcelona  ·  Session 3 Closeout  ·  2026-05-26",
             0.4, 7.05, 9, 0.35, size=9, color=MUTED)
    add_text(slide, f"{page_num} / {total}", 12.2, 7.05, 0.9, 0.35,
             size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 12

# -------------------------------------------------- Slide 1 — Title
s = slide_blank(prs)
add_bg(s, INK)
add_text(s, "MYCORRHIZAL BARCELONA",
         0.7, 1.6, 12, 0.7, size=20, color=AMBER, bold=True)
add_text(s, "Session 3 — Closed",
         0.7, 2.3, 12, 1.5, size=64, color=BONE, bold=True)
add_text(s, "CRISP-DM Phase 3 · Data Preparation Pipeline shipped",
         0.7, 4.0, 12, 0.6, size=22, color=BONE)
add_text(s, "v1.1 → v1.2 · 17 stages · 494 cells × 51 cols · deterministic in 5.2s",
         0.7, 4.7, 12, 0.5, size=16, color=AMBER)
add_text(s, "Group 4 · 2026-05-26",
         0.7, 6.4, 12, 0.5, size=14, color=MUTED)

# -------------------------------------------------- Slide 2 — The arc
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s)
add_text(s, "The Session 3 arc, in one slide",
         0.6, 0.4, 12, 0.7, size=30, color=INK, bold=True)
add_text(s, "Five moves from data to defensible priority map",
         0.6, 1.1, 12, 0.5, size=16, color=MOSS)

beats = [
    ("1.", "Started with vetted Session 2 datasets (Ajuntament trees, FungalRoot, GBIF, Urban Atlas, Landsat LST, Sentinel-2 NDVI, BCN boundary)."),
    ("2.", "Built a 14-stage deterministic CRISP-DM Phase 3 pipeline — `src/clean_data.py` — producing a 495-cell × 40-col `scored_grid` with 4 sub-scores (S1 sealed, S2 LST, S3 NDVI, S4 myco) + 3 weight scenarios."),
    ("3.", "Added v1.1 PRPI (Platanus Replacement Priority Index) — 5th term anchored to the Pla Director de l'Arbrat 2017–2037."),
    ("4.", "Ran a 3-stream parallel deep-research review of PRPI v1.1 → surfaced 5 evidence-based contradictions."),
    ("5.", "Shipped v1.2 — refreshed inventory + peer-reviewed VPA + operational scenario alongside EM-optimistic. Same architecture, tighter priors."),
]
y = 1.85
for num, txt in beats:
    add_text(s, num, 0.6, y, 0.4, 0.5, size=18, color=AMBER, bold=True)
    add_text(s, txt, 1.0, y, 11.7, 1.0, size=14, color=INK)
    y += 1.0

add_footer(s, 2, TOTAL)

# -------------------------------------------------- Slide 3 — Pipeline
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s)
add_text(s, "The pipeline — 17 stages, end-to-end",
         0.6, 0.4, 12, 0.7, size=30, color=INK, bold=True)
add_text(s, "`python src/clean_data.py` · 5.2 s · deterministic · 494 cells × 51 cols",
         0.6, 1.1, 12, 0.5, size=14, color=MOSS)

# Two columns of stages
col1 = [
    "1. Load tree inventory (street + park CSVs)",
    "2. Normalize species names",
    "3. FungalRoot v2.0 lookup + Top-20 override",
    "4. Assign mycorrhizal type per tree",
    "5. Build 400m grid clipped to BCN boundary",
    "6. Spatial join trees → grid cells",
    "7. Per-cell tree stats (incl. n_platanus)",
    "8. GBIF fungal occurrences per cell",
    "9. Zonal stats: sealed / LST / NDVI rasters",
]
col2 = [
    "10. S1–S4 sub-scores",
    "11. PRPI v1.1 (EM-optimistic) — NEW",
    "12. VPA allergenicity + species preference — v1.2 NEW",
    "13. PRPI operational scenario — v1.2 NEW",
    "14. 5-term composite scores (A / B / C)",
    "15. Top-15 priority cells (district-constrained)",
    "16. Intervention classification (5-way)",
    "17. Invariants + write GeoJSON / Parquet",
    "",
]
add_bullets(s, col1, 0.6, 1.8, 6.0, 5.0, size=13)
add_bullets(s, col2, 6.8, 1.8, 6.0, 5.0, size=13)

add_footer(s, 3, TOTAL)

# -------------------------------------------------- Slide 4 — PRPI v1.1
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s)
add_text(s, "PRPI v1.1 — what we shipped first",
         0.6, 0.4, 12, 0.7, size=30, color=INK, bold=True)
add_text(s, "Platanus Replacement Priority Index — anchored to Pla Director 2017–2037",
         0.6, 1.1, 12, 0.5, size=14, color=MOSS)

# Formula card
fbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(1.9), Inches(7), Inches(2.0))
fbox.fill.solid()
fbox.fill.fore_color.rgb = INK
fbox.line.fill.background()
add_text(s, "PRPI formula",
         0.85, 2.0, 6.5, 0.5, size=14, color=AMBER, bold=True)
add_text(s, "PRPI = 0.40 × platanus_pct/100",
         0.85, 2.5, 6.5, 0.4, size=14, color=BONE, font="Consolas")
add_text(s, "     + 0.20 × s3_inverted_ndvi",
         0.85, 2.85, 6.5, 0.4, size=14, color=BONE, font="Consolas")
add_text(s, "     + 0.20 × s4_shift_potential   (EM-optimistic)",
         0.85, 3.20, 6.5, 0.4, size=14, color=BONE, font="Consolas")
add_text(s, "     + 0.20 × (1 − s1_sealed)",
         0.85, 3.55, 6.5, 0.4, size=14, color=BONE, font="Consolas")

# Six locked decisions
add_text(s, "Six locked design decisions",
         8.0, 1.9, 5.0, 0.4, size=16, color=INK, bold=True)
locked = [
    "EM-optimistic replacement (Q. ilex / P. halepensis)",
    "One merged score (PRPI folded into composite)",
    "Both enum + flag (species-replacement + replacement_priority)",
    "Platanus-only counting (n_platanus)",
    "Explicit 2037 Master Plan anchor",
    "s4_shift_ceiling_reached honesty flag",
]
add_bullets(s, locked, 8.0, 2.35, 5.0, 4.0, size=12, color=INK)

add_text(s, "Result: 47 cols · 495 cells · 15 strict-gate cells · 3 species-replacement",
         0.6, 6.3, 12, 0.5, size=14, color=MOSS, bold=True)

add_footer(s, 4, TOTAL)

# -------------------------------------------------- Slide 5 — Deep research kicked off
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s)
add_text(s, "Then we stress-tested it",
         0.6, 0.4, 12, 0.7, size=30, color=INK, bold=True)
add_text(s, "3-stream parallel deep-research review (academia / practice / open source)",
         0.6, 1.1, 12, 0.5, size=14, color=MOSS)

# Three columns
def stream_card(x, title, color, items):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(1.9), Inches(4.0), Inches(5.0))
    card.fill.solid()
    card.fill.fore_color.rgb = BONE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    add_text(s, title, x + 0.2, 2.0, 3.6, 0.5, size=16, color=color, bold=True)
    add_bullets(s, items, x + 0.2, 2.6, 3.6, 4.2, size=11, color=INK)

stream_card(0.6, "Stream 1 · Academia", MOSS, [
    "Platanus pollen allergy in Mediterranean",
    "Urban pollen dispersion models (SILAM / CAMS / COSMO-ART)",
    "Allergenicity rankings of replacement species (VPA / OPALS)",
    "Mycorrhizal-type effects in urban substrate (Verbeek 2025)",
    "Aerobiology monitoring networks (REA / XAC / EAN / AutoPollen)",
])
stream_card(4.8, "Stream 2 · Practice", AMBER, [
    "Pla Director 2017–2037 verification",
    "Eixos Verds + Superilla projects",
    "ASPB / XAC engagement",
    "Comparators: Madrid · Sevilla · Marseille · Rome · Athens",
    "Pilot species performance",
    "Public engagement (Decidim)",
])
stream_card(9.0, "Stream 3 · Open Source", RED, [
    "Pollen models: CAMS · SILAM · HYSPLIT · ENFUSER",
    "Tree inventory + canopy: DeepForest · DetecTree · OpenTrees",
    "Mycorrhizal toolkits: FungalRoot v2 · FUNGuild · GlobalAMFungi",
    "Allergenicity datasets: OPALS · PIA Catalonia · OpenAQ",
    "Reproducible repos surveyed",
])

add_text(s, "Output: 5,800-word APA 7 report · 30+ refs · `outputs/deep-research-platanus-prpi.md`",
         0.6, 7.0, 12, 0.3, size=12, color=MUTED)

add_footer(s, 5, TOTAL)

# -------------------------------------------------- Slide 6 — Five contradictions
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s, color=RED)
add_text(s, "What the evidence said about PRPI v1.1",
         0.6, 0.4, 12, 0.7, size=30, color=INK, bold=True)
add_text(s, "Five contradictions — each tightens v1.1 rather than overturning it",
         0.6, 1.1, 12, 0.5, size=14, color=RED)

findings = [
    ("1.", "The asthma claim doesn't survive scrutiny",
     "Osborne et al. (2017, London n=8.2M): NO significant Platanus → asthma association at any lag. Grass pollen DID. Real burden is rhinoconjunctivitis + Pla a 3 food-allergy cross-reactivity (peach / walnut / hazelnut)."),
    ("2.", "Quercus ilex inverts the public-health goal",
     "Que i 1 = Bet v 1 homolog, VPA IV–V — same allergenicity class as Platanus (Cariñanos & Marinangeli 2021; González-Mancebo et al. 2020). Substitution shifts the peak, doesn't reduce it."),
    ("3.", "AM → EM substrate effect is hypothesis, not delivered outcome",
     "Verbeek et al. (2025, Amsterdam) and Gaimaro et al. (2025, Fairfax VA): urban AM communities shift composition rather than collapse, engineered substrate co-drives colonization."),
    ("4.", "Barcelona already pilots Zelkova + Pistacia, not Q. ilex",
     "Espais Verds operational documentation names Zelkova serrata and Pistacia chinensis as drought-tolerant low-VPA pilots. Q. ilex is not the headline replacement."),
    ("5.", "Inventory snapshot was stale",
     "Brief: 42,828 / 22.6%. Municipal canon (2026): ~43,722 / 27.5%. Open Data BCN `arbrat-viari` 2026_1T release is fresher."),
]
y = 1.85
for num, title, body in findings:
    add_text(s, num, 0.6, y, 0.4, 0.4, size=16, color=RED, bold=True)
    add_text(s, title, 1.0, y, 12, 0.4, size=14, color=INK, bold=True)
    add_text(s, body, 1.0, y + 0.4, 11.7, 0.55, size=11, color=MUTED)
    y += 1.0

add_footer(s, 6, TOTAL)

# -------------------------------------------------- Slide 7 — v1.2 pivot
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s, color=AMBER)
add_text(s, "The v1.2 pivot — small, additive, fully explainable",
         0.6, 0.4, 12, 0.7, size=28, color=INK, bold=True)
add_text(s, "Same architecture · same 2037 anchor · same 6 locked decisions · tighter priors",
         0.6, 1.1, 12, 0.5, size=14, color=AMBER)

changes_keep = [
    "Composite 5-term architecture",
    "Pla Director 2017–2037 policy anchor",
    "Six v1.1 locked design decisions",
    "replacement_priority strict gate",
    "s4_shift_ceiling_reached honesty flag",
    "EM-optimistic scenario (now: upper-bound sensitivity)",
]
changes_add = [
    "Refreshed arbrat-viari → 2026_1T snapshot (188,991 trees)",
    "Cariñanos & Marinangeli 2021 VPA scale (40 species CSV)",
    "cell_vpa_score column (count-weighted mean VPA)",
    "vpa_replacement_delta column",
    "species_preference_present column",
    "prpi_operational scenario (Zelkova/Pistacia palette)",
    "Re-scoped docstrings: asthma → rhinoconjunctivitis + Pla a 3",
    "s4_shift_potential reframed as upper-bound hypothesis",
]

add_text(s, "Kept", 0.6, 1.85, 6.0, 0.4, size=18, color=MOSS, bold=True)
add_bullets(s, changes_keep, 0.6, 2.3, 6.0, 4.5, size=12)

add_text(s, "Added", 6.8, 1.85, 6.0, 0.4, size=18, color=AMBER, bold=True)
add_bullets(s, changes_add, 6.8, 2.3, 6.0, 4.5, size=12)

add_footer(s, 7, TOTAL)

# -------------------------------------------------- Slide 8 — Headline stats
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s)
add_text(s, "v1.2 headline stats (2026_1T snapshot)",
         0.6, 0.4, 12, 0.7, size=30, color=INK, bold=True)
add_text(s, "Pipeline run 2026-05-26 · 5.2 s · all 17 invariants pass",
         0.6, 1.1, 12, 0.5, size=14, color=MOSS)

# Stat cards
def stat_card(x, y, label, value, sublabel, color=MOSS):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(3.8), Inches(1.6))
    card.fill.solid()
    card.fill.fore_color.rgb = BONE
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    add_text(s, label, x + 0.2, y + 0.1, 3.4, 0.3, size=11, color=MUTED)
    add_text(s, value, x + 0.2, y + 0.4, 3.4, 0.8, size=32, color=color, bold=True)
    add_text(s, sublabel, x + 0.2, y + 1.15, 3.4, 0.3, size=10, color=MUTED)

stat_card(0.6, 1.9, "Cells × columns", "494 × 51", "v1.1 was 495 × 47")
stat_card(4.7, 1.9, "Platanus trees in grid", "42,815", "of 42,828 baseline (0.03% boundary edge)")
stat_card(8.8, 1.9, "Pipeline runtime", "5.2 s", "deterministic, single-machine")

stat_card(0.6, 3.7, "prpi range", "[0.15, 0.83]", "EM-optimistic scenario")
stat_card(4.7, 3.7, "prpi_operational", "[0.15, 0.73]", "Zelkova/Pistacia palette · 0.10 lower peak", color=AMBER)
stat_card(8.8, 3.7, "Threshold disagreement", "17 cells", "where the policy choice matters most", color=RED)

stat_card(0.6, 5.5, "replacement_priority", "15", "strict-gate flagged cells")
stat_card(4.7, 5.5, "Ceiling-reached", "164", "AM-blindness honestly flagged")
stat_card(8.8, 5.5, "Intervention enum", "459 / 24 / 7 / 3 / 1", "de-paving · cooling · planting · replace · multi")

add_footer(s, 8, TOTAL)

# -------------------------------------------------- Slide 9 — The disagreement set
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s, color=RED)
add_text(s, "The 17-cell disagreement set",
         0.6, 0.4, 12, 0.7, size=30, color=INK, bold=True)
add_text(s, "Cells where EM-optimistic and operational scenarios disagree at PRPI > 0.5",
         0.6, 1.1, 12, 0.5, size=14, color=RED)

add_text(s, "Why these 17 cells matter most",
         0.6, 1.9, 12, 0.5, size=18, color=INK, bold=True)

reasons = [
    "EM-optimistic (v1.1): 32 cells above threshold — assumes Q. ilex / P. halepensis substitution.",
    "Operational (v1.2): 15 cells above threshold — assumes Zelkova / Pistacia / Sophora substitution.",
    "17 cells are in only one set — the species-choice changes the recommendation here.",
    "These are the cells where Direcció d'Espais Verds should preference the Zelkova/Pistacia palette over a Q. ilex spec.",
    "Both scenarios remain auditable — downstream consumers can compare per-cell.",
]
add_bullets(s, reasons, 0.6, 2.5, 12.2, 3.0, size=14)

add_text(s, "Pivot logic — \"this is what the city is actually doing\"",
         0.6, 5.7, 12, 0.5, size=16, color=AMBER, bold=True)
add_text(s, "Barcelona's Espais Verds pilots Zelkova + Pistacia precisely because they sit in VPA class I–III, are drought-tolerant under +2 °C, and avoid the Platanus → Pla a 3 food-allergy bridge. The EM-optimistic scenario remains useful as a mycorrhizal-ecology upper-bound test; the operational scenario is what gets recommended.",
         0.6, 6.15, 12.2, 0.9, size=11, color=MUTED)

add_footer(s, 9, TOTAL)

# -------------------------------------------------- Slide 10 — Deliverables
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s)
add_text(s, "What ships in the repo",
         0.6, 0.4, 12, 0.7, size=30, color=INK, bold=True)
add_text(s, "Session 3 deliverables — committed, deterministic, reproducible",
         0.6, 1.1, 12, 0.5, size=14, color=MOSS)

# Two columns
code = [
    "src/clean_data.py — 17-stage pipeline",
    "data/raw/vpa-mediterranean-species.csv — 40-species VPA table",
    "data/arbrat-viari.csv — 2026_1T snapshot",
    "data/arbrat-viari-prev-snapshot.csv — diff/rollback",
    "data/processed/scored_grid.geojson + .parquet",
    "requirements.txt — Python 3.11 wheels",
]
docs = [
    "phase-3/data-contract.yaml — schema v1.2.0",
    "docs/data-cleaning-log.md — 16 transforms",
    "docs/pipeline-architecture-v1.md — 17-component spec",
    "outputs/pipeline-results-interpretation.md — §1–§10",
    "outputs/deep-research-platanus-prpi.md — 5,800 words, 30+ refs",
    "HANDOFF.md — Session 4 entry point",
]

add_text(s, "Code & data", 0.6, 1.85, 6.0, 0.4, size=18, color=MOSS, bold=True)
add_bullets(s, code, 0.6, 2.3, 6.0, 4.5, size=12, font="Consolas")

add_text(s, "Documentation", 6.8, 1.85, 6.0, 0.4, size=18, color=AMBER, bold=True)
add_bullets(s, docs, 6.8, 2.3, 6.0, 4.5, size=12, font="Consolas")

add_footer(s, 10, TOTAL)

# -------------------------------------------------- Slide 11 — Limitations + honesty
s = slide_blank(prs)
add_bg(s)
add_accent_bar(s, color=MUTED)
add_text(s, "What this pipeline is NOT for",
         0.6, 0.4, 12, 0.7, size=30, color=INK, bold=True)
add_text(s, "Per Session 4 lecture: every model card states at least 3 NOTs.  We have 5.",
         0.6, 1.1, 12, 0.5, size=14, color=MUTED)

nots = [
    ("NOT a severe-asthma intervention tool",
     "Osborne et al. (2017) found no Platanus → asthma association. Re-scope to rhinoconjunctivitis + Pla a 3 food-allergy."),
    ("NOT a guaranteed mycorrhizal-restoration map",
     "s4_shift_potential is an upper-bound hypothesis — engineered substrate may dominate over host species effects (Verbeek et al. 2025)."),
    ("NOT a high-resolution pollen-dispersion forecast",
     "No European model resolves the 400m grid (SILAM/CAMS at ~10 km). Local point ground-truth via XAC / PIA UAB only."),
    ("NOT a substitute for clinical or regulatory air-quality decisions",
     "VPA proxies allergenic potential, not exposure or sensitization in any specific cohort."),
    ("NOT valid outside the trained spatial extent",
     "Index is bounded to Barcelona's 400m grid. Peri-urban (Collserola / Garraf) reference patch is still missing — flagged in Session 4 agenda."),
]
y = 1.95
for title, body in nots:
    add_text(s, "✕", 0.6, y, 0.4, 0.5, size=18, color=RED, bold=True)
    add_text(s, title, 1.0, y, 12, 0.4, size=14, color=INK, bold=True)
    add_text(s, body, 1.0, y + 0.4, 11.7, 0.55, size=11, color=MUTED)
    y += 1.0

add_footer(s, 11, TOTAL)

# -------------------------------------------------- Slide 12 — Session 4
s = slide_blank(prs)
add_bg(s, INK)
add_text(s, "Session 4 — Modeling Phase",
         0.7, 0.7, 12, 0.7, size=20, color=AMBER, bold=True)
add_text(s, "scored_grid.parquet is the modeling-ready input",
         0.7, 1.4, 12, 0.8, size=36, color=BONE, bold=True)

# Four boxes for the four lecture themes
def s4_card(x, y, title, body):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(5.9), Inches(2.0))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x25, 0x3A, 0x2C)
    card.line.color.rgb = AMBER
    card.line.width = Pt(1.5)
    add_text(s, title, x + 0.2, y + 0.15, 5.5, 0.4, size=16, color=AMBER, bold=True)
    add_text(s, body, x + 0.2, y + 0.6, 5.5, 1.3, size=12, color=BONE)

s4_card(0.7, 2.7, "Split strategy",
        "Spatial / clustered hold-out (NOT random rows).  Hold out districts or k-means on (lat, lon).  Test set is locked.")
s4_card(6.8, 2.7, "Baselines",
        "Dumb mean · spatial nearest · domain heuristic (e.g. Eixample = high barrier).  Beating these is the bar.")
s4_card(0.7, 4.9, "Metrics + model card",
        "Linear models first.  Report ≥ 3 NOTs.  Train / eval / test metrics — never just one number.")
s4_card(6.8, 4.9, "Sensitivity scenarios",
        "Both prpi (EM-optimistic) and prpi_operational available.  Scenario B composite_score_B is the primary target.")

add_text(s, "Read HANDOFF.md to resume.  CRISP-DM Phase 4 is loaded and ready.",
         0.7, 7.0, 12, 0.4, size=12, color=AMBER)

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
